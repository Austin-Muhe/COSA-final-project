from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
REPORT = ROOT / "report"
OUT_PPTX = ROOT / "COSA_Dynamic_Regime_Presentation.pptx"
OUT_PDF = ROOT / "COSA_Dynamic_Regime_Presentation.pdf"

ASSETS = {
    "clean_ablation": ROOT / "results/final_figures/figure2_etth1_ablation_bar.png",
    "recovery": ROOT / "results/final_figures/figure1_etth1_h720_blackswan_recovery.png",
    "weather_tsne": ROOT / "results/regime_adapter_latents_mixed_loss/tsne/DLinear/weather/96/latent_regimes.png",
    "etth1_tsne": ROOT / "results/regime_adapter_latents_etth1_720/tsne/DLinear/ETTh1/720/latent_regimes.png",
}

SLIDE_W, SLIDE_H = 13.333, 7.5
COL = {
    "ink": RGBColor(15, 23, 42),
    "muted": RGBColor(71, 85, 105),
    "line": RGBColor(203, 213, 225),
    "blue": RGBColor(37, 99, 235),
    "sky": RGBColor(224, 242, 254),
    "green": RGBColor(22, 163, 74),
    "green_bg": RGBColor(220, 252, 231),
    "amber": RGBColor(180, 83, 9),
    "amber_bg": RGBColor(254, 243, 199),
    "red": RGBColor(220, 38, 38),
    "red_bg": RGBColor(254, 226, 226),
    "purple": RGBColor(124, 58, 237),
    "purple_bg": RGBColor(237, 233, 254),
    "white": RGBColor(255, 255, 255),
    "bg": RGBColor(248, 250, 252),
}


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COL["bg"]


def add_text(slide, text, x, y, w, h, size=24, bold=False, color="ink", align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = COL[color]
    return box


def add_bullets(slide, items, x, y, w, h, size=22, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = COL[color]
        p.space_after = Pt(8)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.28, 11.9, 0.48, 28, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.82, 11.8, 0.34, 14, False, "muted")
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.55), Inches(1.18), Inches(12.75), Inches(1.18))
    line.line.color.rgb = COL["line"]
    line.line.width = Pt(1.1)


def add_box(slide, text, x, y, w, h, fill="white", outline="line", size=20, bold=True, color="ink"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COL[fill]
    shape.line.color.rgb = COL[outline]
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COL[color]
    return shape


def add_arrow(slide, x1, y1, x2, y2, color="muted", width=2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COL[color]
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_image(slide, path: Path, x, y, w=None, h=None):
    if w is None and h is None:
        raise ValueError("set w or h")
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    im = Image.open(path)
    iw, ih = im.size
    aspect = iw / ih
    if w is not None:
        h = w / aspect
    else:
        w = h * aspect
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = text


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_text(s, "COSA+ and Dynamic Regime Adaptation", 0.7, 0.55, 11.8, 0.65, 34, True)
    add_text(s, "Time-Series Test-Time Adaptation under Distribution Shift", 0.72, 1.25, 11.5, 0.35, 19, False, "muted")
    add_box(s, "Research Question", 0.9, 2.05, 3.0, 0.65, "sky", "blue", 21, True, "blue")
    add_text(s, "Can output-space adaptation stay useful when deployed time-series streams suddenly change?", 1.0, 2.95, 11.1, 0.7, 27, True, "ink", "center")
    add_box(s, "COSA+\nstronger online correction", 1.15, 4.45, 3.5, 1.15, "green_bg", "green", 21, True, "green")
    add_box(s, "Dynamic Regime Adapter\nzero-shot, no backprop", 4.9, 4.45, 3.7, 1.15, "purple_bg", "purple", 21, True, "purple")
    add_box(s, "Compare trade-offs\nnot a universal winner", 8.85, 4.45, 3.35, 1.15, "amber_bg", "amber", 21, True, "amber")
    add_text(s, "CPS3830 Course Project | 7-minute presentation", 0.78, 6.72, 11.8, 0.25, 12, False, "muted")
    add_notes(s, "Start with the research question. Our project is about test-time adaptation for time-series forecasting. We compare two directions beyond COSA: COSA+ keeps online adaptation and Dynamic Regime Adapter moves adaptation offline. The key message is trade-offs, not one universal best method.")

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Motivation: Forecasting Streams Do Not Stay Clean", "Clean train/test splits miss abrupt real-world regime changes.")
    add_box(s, "Training regime", 0.9, 2.0, 2.4, 0.85, "sky", "blue", 20, True, "blue")
    add_arrow(s, 3.45, 2.42, 4.7, 2.42)
    add_box(s, "Deployed stream", 4.85, 2.0, 2.4, 0.85, "white", "line", 20)
    add_arrow(s, 7.4, 2.42, 8.65, 2.42, "red")
    add_box(s, "Sudden shift", 8.8, 2.0, 2.4, 0.85, "red_bg", "red", 20, True, "red")
    add_bullets(s, ["Sensor faults", "Weather anomalies", "Demand spikes", "Structural changes"], 1.0, 3.45, 4.4, 2.0, 25)
    add_box(s, "Problem", 6.1, 3.3, 5.8, 0.7, "amber_bg", "amber", 22, True, "amber")
    add_text(s, "A model can look accurate on a clean test split, but fail exactly when the stream enters a new regime.", 6.25, 4.2, 5.5, 1.0, 25, True, "ink", "center")
    add_notes(s, "This slide motivates the problem. In deployment, time series can change because of sensors, weather, demand, or structural changes. Standard evaluation assumes the test distribution is stable. Our project asks what happens when that assumption breaks.")

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Baseline: COSA Online Output-Space Adaptation", "Frozen backbone + lightweight adapter optimized during test time.")
    add_box(s, "Context\nCt", 0.9, 2.05, 1.65, 0.8, "sky", "blue", 20, True, "blue")
    add_box(s, "Frozen\nDLinear", 0.9, 3.45, 1.65, 0.8, "white", "line", 20)
    add_box(s, "Base forecast\nY0", 3.15, 2.75, 1.8, 0.8, "white", "line", 20)
    add_box(s, "Output adapter\nH = W[Y0 || Ct] + b", 5.55, 2.75, 2.55, 0.8, "amber_bg", "amber", 18, True, "amber")
    add_box(s, "Adapted forecast\nY = Y0 + tanh(g)H", 8.85, 2.75, 3.1, 0.8, "green_bg", "green", 18, True, "green")
    add_arrow(s, 2.55, 2.45, 3.15, 3.0)
    add_arrow(s, 2.55, 3.85, 3.15, 3.2)
    add_arrow(s, 4.95, 3.15, 5.55, 3.15)
    add_arrow(s, 8.1, 3.15, 8.85, 3.15)
    add_box(s, "Strength", 1.15, 5.1, 2.1, 0.45, "green_bg", "green", 17, True, "green")
    add_text(s, "Small module; no full model retraining", 0.9, 5.65, 3.0, 0.6, 20, True, "ink", "center")
    add_box(s, "Cost", 5.45, 5.1, 2.1, 0.45, "red_bg", "red", 17, True, "red")
    add_text(s, "Still needs online gradient updates", 5.1, 5.65, 2.9, 0.6, 20, True, "ink", "center")
    add_box(s, "Question", 9.25, 5.1, 2.1, 0.45, "amber_bg", "amber", 17, True, "amber")
    add_text(s, "What helps under abrupt shifts?", 8.9, 5.65, 2.9, 0.6, 20, True, "ink", "center")
    add_notes(s, "COSA is our baseline. It freezes the forecasting backbone and adapts a small output-space module online. This is lighter than updating the full model, but it still requires test-time backpropagation.")

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Our Two Directions", "One stays online; one moves adaptation offline.")
    add_box(s, "COSA+", 0.8, 1.55, 2.3, 0.65, "green_bg", "green", 24, True, "green")
    add_box(s, "Dynamic Regime Adapter", 7.05, 1.55, 3.6, 0.65, "purple_bg", "purple", 24, True, "purple")
    add_box(s, "Mean + std\ncontext", 0.85, 2.85, 1.9, 0.75, "white", "line", 18)
    add_box(s, "Vector gate\ng ∈ R^T", 3.15, 2.85, 1.9, 0.75, "white", "line", 18)
    add_box(s, "Step-wise\ncorrection", 5.45, 2.85, 1.9, 0.75, "green_bg", "green", 18, True, "green")
    add_arrow(s, 2.75, 3.22, 3.15, 3.22, "green")
    add_arrow(s, 5.05, 3.22, 5.45, 3.22, "green")
    add_box(s, "Context X", 7.0, 2.65, 1.5, 0.62, "sky", "blue", 17, True, "blue")
    add_box(s, "RegimeEncoder\nz", 8.85, 2.45, 1.85, 1.0, "white", "line", 17)
    add_box(s, "GateGenerator\ng", 11.0, 2.45, 1.65, 1.0, "purple_bg", "purple", 17, True, "purple")
    add_arrow(s, 8.5, 2.95, 8.85, 2.95, "purple")
    add_arrow(s, 10.7, 2.95, 11.0, 2.95, "purple")
    add_text(s, "Online: adapts during evaluation", 1.05, 4.55, 5.4, 0.45, 23, True, "green", "center")
    add_text(s, "Zero-shot: single forward pass at test time", 7.1, 4.55, 5.3, 0.45, 23, True, "purple", "center")
    add_box(s, "Trade-off axis: robustness behavior vs. inference cost", 2.35, 6.0, 8.6, 0.65, "amber_bg", "amber", 24, True, "amber")
    add_notes(s, "We explored two directions. COSA+ keeps COSA’s online adaptation but adds richer context and a horizon-wise vector gate. Dynamic Regime Adapter learns a latent regime representation offline, then generates a correction gate in one forward pass.")

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Core Experimental Design", "Clean forecasting + black-swan robustness tests.")
    add_box(s, "Datasets\nWeather, ETTh1", 0.75, 1.75, 2.4, 1.0, "sky", "blue", 21, True, "blue")
    add_box(s, "Backbone\nFrozen DLinear", 3.65, 1.75, 2.4, 1.0, "white", "line", 21)
    add_box(s, "Horizons\n96, 720", 6.55, 1.75, 2.4, 1.0, "white", "line", 21)
    add_box(s, "Metric\nMSE / MAE", 9.45, 1.75, 2.4, 1.0, "green_bg", "green", 21, True, "green")
    add_text(s, "Black-Swan shift protocol", 0.85, 3.35, 4.0, 0.45, 24, True)
    add_box(s, "Original stream", 0.95, 4.25, 2.0, 0.65, "white", "line", 18)
    add_arrow(s, 3.1, 4.58, 4.3, 4.58)
    add_box(s, "Shift point\nTshift", 4.45, 4.15, 1.55, 0.85, "red_bg", "red", 18, True, "red")
    add_arrow(s, 6.15, 4.58, 7.35, 4.58)
    add_box(s, "Shifted stream\nx't = S(xt; α)", 7.5, 4.25, 2.3, 0.65, "amber_bg", "amber", 18, True, "amber")
    add_box(s, "Evaluate\nmse_after_50", 10.25, 4.25, 2.0, 0.65, "green_bg", "green", 18, True, "green")
    for i, label in enumerate(["Level", "Variance", "Trend", "Spike"]):
        add_box(s, label, 1.05 + i * 2.25, 5.65, 1.75, 0.5, "purple_bg", "purple", 17, True, "purple")
    add_notes(s, "The experiments use Weather and ETTh1 with a frozen DLinear backbone. We evaluate clean test performance and black-swan shifts. The black-swan protocol transforms the stream after a shift point using level, variance, trend, or spike shifts, then measures recovery with mse_after_50.")

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Main Results", "Clean-test gains are similar; abrupt-shift behavior differs.")
    add_box(s, "Weather / h=96\nBest clean MSE", 0.75, 1.55, 2.5, 0.85, "sky", "blue", 19, True, "blue")
    add_box(s, "Dynamic mixed loss\n+2.48% MSE", 3.45, 1.55, 2.45, 0.85, "purple_bg", "purple", 19, True, "purple")
    add_box(s, "ETTh1 / h=720\nBest clean MSE", 6.25, 1.55, 2.5, 0.85, "sky", "blue", 19, True, "blue")
    add_box(s, "COSA+\n+4.86% MSE", 8.95, 1.55, 2.45, 0.85, "green_bg", "green", 19, True, "green")
    add_image(s, ASSETS["recovery"], 0.75, 2.75, w=5.95)
    add_text(s, "ETTh1 / DLinear / h=720 black-swan recovery", 0.95, 6.15, 5.4, 0.25, 14, True, "muted", "center")
    add_image(s, ASSETS["clean_ablation"], 7.05, 2.78, w=5.45)
    add_text(s, "Ablation: vector gate helps most at long horizons", 7.2, 6.15, 5.1, 0.25, 14, True, "muted", "center")
    add_notes(s, "The main results show the trade-off. On Weather horizon 96, the Dynamic Regime Adapter with mixed loss gives the best MSE without test-time backpropagation. On ETTh1 horizon 720, COSA+ gives the strongest clean and black-swan correction. The recovery plot shows COSA+ is consistently below original COSA after abrupt shifts.")

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Takeaways, Limitations, Future Work", "The answer is conditional, which is the useful result.")
    add_box(s, "What worked", 0.85, 1.55, 2.45, 0.55, "green_bg", "green", 21, True, "green")
    add_bullets(s, ["COSA+ improves long-horizon post-shift recovery", "Dynamic adapter gives zero-shot adaptation", "Vector gate explains most COSA+ gains"], 0.9, 2.25, 4.0, 1.95, 20)
    add_box(s, "Limitations", 5.0, 1.55, 2.45, 0.55, "red_bg", "red", 21, True, "red")
    add_bullets(s, ["Only Weather and ETTh1", "Black-swan shifts are synthetic", "Weather h=720 shows vector gate can hurt", "Dynamic black-swan test still open"], 5.05, 2.25, 3.85, 2.35, 20)
    add_box(s, "Next steps", 9.1, 1.55, 2.45, 0.55, "amber_bg", "amber", 21, True, "amber")
    add_bullets(s, ["Gate regularization", "Shift detection", "Huber / robust losses", "More datasets and horizons"], 9.15, 2.25, 3.45, 2.0, 20)
    add_box(s, "Final message: COSA+ is stronger online under abrupt long-horizon shifts; Dynamic Regime Adapter is cheaper at inference with no test-time backpropagation.", 1.1, 5.55, 11.0, 0.9, "purple_bg", "purple", 22, True, "purple")
    add_notes(s, "Close with what we learned. COSA+ is useful in a specific robustness regime, especially long-horizon abrupt shifts. Dynamic Regime Adapter is attractive when inference cost matters. We do not claim a universal best method; the project compares trade-offs and identifies limitations.")

    prs.save(OUT_PPTX)


def draw_pdf_slide(c, title, bullets, footer=None):
    w, h = landscape((960, 540))
    c.setFillColorRGB(248 / 255, 250 / 255, 252 / 255)
    c.rect(0, 0, w, h, stroke=0, fill=1)
    c.setFillColorRGB(15 / 255, 23 / 255, 42 / 255)
    c.setFont("Helvetica-Bold", 25)
    c.drawString(42, 492, title)
    c.setStrokeColorRGB(203 / 255, 213 / 255, 225 / 255)
    c.line(42, 472, 918, 472)
    c.setFont("Helvetica", 18)
    y = 405
    for b in bullets:
        c.setFillColorRGB(15 / 255, 23 / 255, 42 / 255)
        c.circle(60, y + 6, 4, fill=1, stroke=0)
        c.drawString(78, y, b)
        y -= 46
    if footer:
        c.setFillColorRGB(71 / 255, 85 / 255, 105 / 255)
        c.setFont("Helvetica", 12)
        c.drawString(42, 32, footer)
    c.showPage()


def build_pdf_summary():
    c = canvas.Canvas(str(OUT_PDF), pagesize=landscape((960, 540)))
    slides = [
        ("COSA+ and Dynamic Regime Adaptation", ["Research question: can adaptation stay useful under sudden time-series shifts?", "Compare COSA+ online correction with Dynamic zero-shot adaptation.", "Main claim: trade-offs, not a universal winner."]),
        ("Motivation", ["Deployed streams shift due to sensor faults, weather anomalies, and demand spikes.", "Clean train/test splits hide this failure mode.", "We evaluate both clean performance and abrupt-shift recovery."]),
        ("Baseline: COSA", ["Frozen DLinear backbone.", "Online output-space adapter uses recent context.", "Lightweight, but still requires test-time gradient updates."]),
        ("Our Two Directions", ["COSA+: mean+std context and horizon-wise vector gate.", "Dynamic Regime Adapter: RegimeEncoder z plus GateGenerator g.", "COSA+ adapts online; Dynamic adapts by one forward pass."]),
        ("Experimental Design", ["Datasets: Weather and ETTh1.", "Backbone: frozen DLinear; horizons 96 and 720.", "Black-swan shifts: level, variance, trend, spike; metric: mse_after_50."]),
        ("Main Results", ["Weather h=96: Dynamic mixed loss has best clean MSE (+2.48%).", "ETTh1 h=720: COSA+ has best clean MSE (+4.86%).", "Under ETTh1 h=720 shifts, COSA+ gives strongest post-shift recovery."]),
        ("Takeaways", ["COSA+ helps when long-horizon errors accumulate after abrupt shifts.", "Dynamic Regime Adapter removes test-time backpropagation.", "Future work: gate regularization, shift detection, robust losses, more datasets."]),
    ]
    for title, bullets in slides:
        draw_pdf_slide(c, title, bullets, "PDF summary export of COSA_Dynamic_Regime_Presentation.pptx")
    c.save()


if __name__ == "__main__":
    for name, path in ASSETS.items():
        if not path.exists():
            raise FileNotFoundError(f"Missing asset {name}: {path}")
    build_pptx()
    build_pdf_summary()
    print(f"saved {OUT_PPTX}")
    print(f"saved {OUT_PDF}")
